from flask import Flask, request, jsonify
from flask_cors import CORS
import os
os.environ["CUDA_VISIBLE_DEVICES"] = "-1"  # Disable GPU and force CPU usage
os.environ["TF_CPP_MIN_LOG_LEVEL"] = "2"  # Suppress warnings

import sqlite3
import pdfplumber  # PDF text extraction
from docx import Document  # Word document handling
from pptx import Presentation  # PowerPoint handling
import logging  # Logging module
from transformers import pipeline

app = Flask(__name__)
# CORS(app)
CORS(app, resources={r"/*": {"origins": "*"}}, supports_credentials=True)


# Configure Logging
logging.basicConfig(level=logging.INFO)

# Load Summarization Model
summarizer = pipeline("summarization", model="sshleifer/distilbart-cnn-6-6", device=-1)

# Database Setup
db_path = "summaries.db"
if not os.path.exists(db_path):
    conn = sqlite3.connect(db_path)
    c = conn.cursor()
    c.execute('''CREATE TABLE summaries (id INTEGER PRIMARY KEY, title TEXT, original_text TEXT, summary TEXT, date TEXT)''')
    conn.commit()
    conn.close()

# Function to extract text from PDF
def extract_text_from_pdf(pdf_file):
    with pdfplumber.open(pdf_file) as pdf:
        text = ""
        for page in pdf.pages:
            extracted_text = page.extract_text()
            if extracted_text:
                text += extracted_text + "\n"
    return text.strip()

# Function to extract text from DOCX (Word)
def extract_text_from_docx(docx_file):
    doc = Document(docx_file)
    text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    return text.strip()

# Function to extract text from PPTX (PowerPoint)
def extract_text_from_pptx(pptx_file):
    ppt = Presentation(pptx_file)
    text = ""
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

# Upload Notes & Generate Summary
@app.route('/upload', methods=['POST'])
def upload_notes():
    title = request.form.get("title")
    notes = request.form.get("notes")
    uploaded_file = request.files.get("file")

    if uploaded_file:
        file_ext = uploaded_file.filename.lower().split(".")[-1]
        
        if file_ext == "txt":
            notes = uploaded_file.read().decode("utf-8")
        elif file_ext == "pdf":
            notes = extract_text_from_pdf(uploaded_file)
        elif file_ext == "docx":
            notes = extract_text_from_docx(uploaded_file)
        elif file_ext == "pptx":
            notes = extract_text_from_pptx(uploaded_file)
        else:
            logging.warning("Unsupported file type uploaded: %s", uploaded_file.filename)
            return jsonify({"error": "Unsupported file type"}), 400

    if not notes:
        logging.warning("No text provided for summarization.")
        return jsonify({"error": "No text provided"}), 400

    # Dynamically set max_length based on input size
    input_length = len(notes.split())  # Count words
    max_summary_length = min(150, max(10, input_length // 2))  # Summary should be half of input length, max 150 words

    summary = summarizer(notes, max_length=max_summary_length, min_length=5, do_sample=False)[0]['summary_text']
    
    conn = sqlite3.connect(db_path)
    c = conn.cursor()
    c.execute("INSERT INTO summaries (title, original_text, summary, date) VALUES (?, ?, ?, datetime('now'))", (title, notes, summary))
    conn.commit()
    conn.close()

    logging.info("Summary successfully generated for: %s", title)

    return jsonify({"message": "Summary generated successfully", "summary": summary})

# Retrieve All Summaries
@app.route('/summaries', methods=['GET'])
def get_summaries():
    conn = sqlite3.connect(db_path)
    c = conn.cursor()
    c.execute("SELECT id, title, summary, date FROM summaries")
    summaries = c.fetchall()
    conn.close()
    
    # Convert results to a JSON-friendly format
    result = [
        {"id": row[0], "title": row[1], "summary": row[2], "date": row[3]}
        for row in summaries
    ]
    
    return jsonify(result)

# Delete Summary by ID
@app.route('/summary/<int:summary_id>', methods=['OPTIONS', 'DELETE'])
def delete_summary(summary_id):
    if request.method == 'OPTIONS':
        # Preflight request - reply with headers only
        response = jsonify({"message": "Preflight request successful"})
        response.headers.add("Access-Control-Allow-Origin", "*")
        response.headers.add("Access-Control-Allow-Methods", "DELETE, OPTIONS")
        response.headers.add("Access-Control-Allow-Headers", "Content-Type, Authorization")
        return response, 200

    conn = sqlite3.connect(db_path)
    c = conn.cursor()
    c.execute("DELETE FROM summaries WHERE id = ?", (summary_id,))
    conn.commit()
    conn.close()

    response = jsonify({"message": "Summary deleted successfully", "id": summary_id})
    response.headers.add("Access-Control-Allow-Origin", "*")
    return response, 200



if __name__ == '__main__':
    app.run(debug=False)
